# -*- coding: utf-8 -*-
"""通用课堂教学 PPT 生成器（数据驱动，courseware-pptx 标准版 v2）。

接受大纲 JSON，按 kind 分发渲染，继承模板（16:9/主题/母版）、套用户锁定规范：
标题微软雅黑加粗居中36pt、正文黑体加粗≥24pt、数学符号 Times 斜体、红标 #FF0000、
行距1.5、视频按 kind 嵌入。

用法：
  python gen_courseware.py 大纲.json --out 课件.pptx --template 课件模板.pptx
"""
import os
import re
import json
import argparse
from lxml import etree
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ---- 锁定配色 ----
BLUE = RGBColor(0x4F, 0x81, 0xBD)
RED = RGBColor(0xFF, 0x00, 0x00)
DARK = RGBColor(0x1F, 0x4F, 0x7D)
BLACK = RGBColor(0x2E, 0x2E, 0x2E)

W = Inches(13.333)
H = Inches(7.5)
MARGIN_LEFT = Inches(0.7)
MARGIN_RIGHT = Inches(0.7)
MARGIN_TOP = Inches(0.55)
BODY_TOP = Inches(1.5)
BODY_WIDTH = W - MARGIN_LEFT - MARGIN_RIGHT


# ---------------------------------------------------------------------------
# 字体助手
# ---------------------------------------------------------------------------
def set_run(run, text, *, size=24, bold=False, italic=False, color=None,
            latin='微软雅黑', ea='微软雅黑'):
    run.text = text
    f = run.font
    if size:
        f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    if color is not None:
        f.color.rgb = color
    f.name = latin
    rPr = run._r.get_or_add_rPr()
    ns_a = qn('a:ea')
    ea_el = rPr.find(ns_a)
    if ea_el is None:
        ea_el = etree.SubElement(rPr, ns_a)
    ea_el.set('typeface', ea)
    return run


def add_marked_runs(paragraph, text, *, size=24, bold=False, italic=False, color=BLACK,
                    latin='微软雅黑', ea='微软雅黑'):
    """标记：数学 $..$ 或 $$..$$（Times斜体）、**加粗**、!!红字!!"""
    parts = re.split(r'(\$\$[^$]*\$\$|\$[^$]+\$|\*\*.*?\*\*|!!.*?!!)', text)
    for part in parts:
        if not part:
            continue
        if part.startswith('$$') and part.endswith('$$'):
            set_run(paragraph.add_run(), part[2:-2], size=size, bold=False, italic=True,
                    color=color, latin='Times New Roman', ea='Times New Roman')
        elif part.startswith('$') and part.endswith('$'):
            set_run(paragraph.add_run(), part[1:-1], size=size, bold=False, italic=True,
                    color=color, latin='Times New Roman', ea='Times New Roman')
        elif part.startswith('**') and part.endswith('**'):
            set_run(paragraph.add_run(), part[2:-2], size=size, bold=True, italic=italic,
                    color=color, latin=latin, ea=ea)
        elif part.startswith('!!') and part.endswith('!!'):
            set_run(paragraph.add_run(), part[2:-2], size=size, bold=bold, italic=italic,
                    color=RED, latin=latin, ea=ea)
        else:
            set_run(paragraph.add_run(), part, size=size, bold=bold, italic=italic,
                    color=color, latin=latin, ea=ea)
    paragraph.line_spacing = 1.5
    paragraph.space_after = Pt(10)


# ---------------------------------------------------------------------------
# 布局助手
# ---------------------------------------------------------------------------
def delete_all_slides(prs):
    while len(prs.slides) > 0:
        rId = prs.slides._sldIdLst[0].rId
        prs.part.drop_rel(rId)
        prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])


def find_blank_layout(prs):
    for layout in prs.slide_layouts:
        if '空白' in layout.name or 'Blank' in layout.name:
            return layout
    best = prs.slide_layouts[-1]
    for layout in prs.slide_layouts:
        if len(layout.placeholders) <= len(best.placeholders):
            best = layout
    return best


def add_title_tag(slide, text):
    """顶部红色【标签】标题（居中）"""
    tx = slide.shapes.add_textbox(MARGIN_LEFT, MARGIN_TOP, Inches(11.9), Inches(0.9))
    tf = tx.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    add_marked_runs(p, text, size=36, bold=True, color=RED)
    return tx


def add_body_box(slide, bullets, *, top=BODY_TOP, left=MARGIN_LEFT,
                 width=BODY_WIDTH, height=None, size=26, ea='黑体', bold=True):
    """内容区：黑体加粗正文，每段一个 bullet"""
    if height is None:
        height = Inches(5.8) - (top - MARGIN_TOP)
    tx = slide.shapes.add_textbox(left, top, width, height)
    tf = tx.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.5
        p.space_after = Pt(12)
        add_marked_runs(p, '· ' + b, size=size, bold=bold, latin=ea, ea=ea)
    return tx


def add_sub_tag(slide, text, *, top=Inches(1.05), color=DARK):
    tx = slide.shapes.add_textbox(MARGIN_LEFT, top, Inches(11.9), Inches(0.45))
    tf = tx.text_frame
    p = tf.paragraphs[0]
    add_marked_runs(p, text, size=26, bold=False, color=color)
    return tx


def embed_video(slide, path, left, top, w, h, poster=None):
    if not path or not os.path.exists(path):
        return False
    try:
        slide.shapes.add_movie(
            path, left, top, w, h,
            poster_frame_image=poster if (poster and os.path.exists(poster)) else None,
            mime_type='video/mp4')
        return True
    except Exception as e:
        print(f"视频嵌入警告：{e}")
        return False


# ---------------------------------------------------------------------------
# 主构建
# ---------------------------------------------------------------------------
def build(data, out, template=None, cover=None):
    if template and os.path.exists(template):
        prs = Presentation(template)
        if abs(prs.slide_width / 914400 - 13.333) > 0.1:
            prs.slide_width = W
            prs.slide_height = H
    else:
        prs = Presentation()
        prs.slide_width = W
        prs.slide_height = H
    delete_all_slides(prs)
    blank = find_blank_layout(prs)

    for sd in data['pages']:
        slide = prs.slides.add_slide(blank)
        kind = sd.get('kind', 'section')

        if kind == 'cover':
            # 章节封面模式：主标题显示"八年级 第X章 章名"，副标题保留课时标题
            if cover and cover.get('grade') and cover.get('chapter') and cover.get('name'):
                cover_title = "%s　第%s章　%s" % (cover['grade'], cover['chapter'], cover['name'])
                cover_sub = sd.get('title', '') or sd.get('subtitle', '')
            else:
                cover_title = sd.get('title', '') or ''
                cover_sub = sd.get('subtitle', '')
            title = slide.shapes.add_textbox(Inches(1), Inches(2.3), Inches(11.3), Inches(1.4))
            tf = title.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            add_marked_runs(p, cover_title, size=54, bold=True, color=DARK)
            sub = slide.shapes.add_textbox(Inches(1), Inches(4.0), Inches(11.3), Inches(0.8))
            p2 = sub.text_frame.paragraphs[0]
            p2.alignment = PP_ALIGN.CENTER
            add_marked_runs(p2, cover_sub, size=32, bold=False, color=BLUE)
            info = slide.shapes.add_textbox(Inches(1), Inches(5.0), Inches(11.3), Inches(0.6))
            p3 = info.text_frame.paragraphs[0]
            p3.alignment = PP_ALIGN.CENTER
            add_marked_runs(p3, sd.get('info', ''), size=24, bold=False, color=BLACK)

        elif kind == 'section':
            add_title_tag(slide, sd['tag'])
            add_body_box(slide, sd.get('bullets', []), size=sd.get('size', 26))

        elif kind == 'import':
            add_title_tag(slide, sd['tag'])
            add_body_box(slide, sd.get('bullets', []), size=sd.get('size', 26))
            if sd.get('video_path'):
                embed_video(slide, sd['video_path'], Inches(3.15), Inches(3.55),
                            Inches(7.0), Inches(3.94), poster=sd.get('poster_path'))

        elif kind == 'study':
            add_title_tag(slide, sd['tag'])
            top = BODY_TOP
            if sd.get('sub_tag'):
                add_sub_tag(slide, sd['sub_tag'], top=Inches(1.05))
                top = Inches(1.55)
            if sd.get('video_path'):
                # 左视频 + 右文字（与 stability 同版式）
                embed_video(slide, sd['video_path'], Inches(0.6), Inches(1.65),
                            Inches(6.3), Inches(3.78), poster=sd.get('poster_path'))
                add_body_box(slide, sd.get('bullets', []), top=Inches(1.75),
                             left=Inches(7.1), width=Inches(5.55), size=24)
            else:
                body_height = Inches(5.9) - (top - MARGIN_TOP)
                add_body_box(slide, sd.get('bullets', []), top=top, height=body_height,
                             size=sd.get('size', 26))

        elif kind == 'practice':
            add_title_tag(slide, sd['tag'])
            add_body_box(slide, sd.get('bullets', []), top=BODY_TOP, size=sd.get('size', 25))

        elif kind == 'stability':
            add_title_tag(slide, sd['tag'])
            add_sub_tag(slide, '实验演示', top=Inches(1.05))
            if sd.get('video_path'):
                embed_video(slide, sd['video_path'], Inches(0.6), Inches(1.65),
                            Inches(6.3), Inches(3.78))
            add_body_box(slide, sd.get('bullets', []), top=Inches(1.75),
                         left=Inches(7.1), width=Inches(5.55), size=26)

        elif kind == 'summary':
            add_title_tag(slide, sd['tag'])
            add_body_box(slide, sd.get('bullets', []), top=BODY_TOP, size=sd.get('size', 27))

        else:
            add_title_tag(slide, sd.get('tag', sd.get('title', '')))
            add_body_box(slide, sd.get('bullets', []), size=sd.get('size', 26))

    prs.save(out)
    print("OK ->", out, "页数:", len(data['pages']))


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument('json')
    ap.add_argument('--out', default='课件.pptx')
    ap.add_argument('--template', default=None)
    ap.add_argument('--grade', default=None, help='章节封面年级，如 八年级')
    ap.add_argument('--chapter', default=None, help='章节封面章号，如 13')
    ap.add_argument('--chapter-name', default=None, help='章节封面章名，如 三角形')
    args = ap.parse_args()
    with open(args.json, encoding='utf-8') as f:
        data = json.load(f)
    cover = None
    if args.grade and args.chapter and args.chapter_name:
        cover = {'grade': args.grade, 'chapter': args.chapter,
                 'name': args.chapter_name}
    build(data, args.out, args.template, cover)


if __name__ == '__main__':
    main()
